from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── color palette ──
BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
ACCENT2 = RGBColor(0x00, 0xE5, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xB0, 0xC4, 0xDE)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xFF, 0x55, 0x55)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
CARD_BG = RGBColor(0x25, 0x25, 0x45)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, title, body, title_color=ACCENT2, body_color=LIGHT):
    card = add_shape(slide, left, top, width, height, fill_color=CARD_BG)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=18, bold=True, color=title_color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(0.9),
                 body, font_size=14, color=body_color)
    return card

def add_arrow(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()
    return arrow

def add_stage_box(slide, left, top, width, height, label, sub_label, color=ACCENT):
    shape = add_shape(slide, left, top, width, height, fill_color=color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    tf.add_paragraph().text = sub_label
    tf.paragraphs[1].font.size = Pt(11)
    tf.paragraphs[1].font.color.rgb = WHITE
    tf.paragraphs[1].font.name = "Arial"
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    return shape

# ════════════════════════════════════════════════
# Slide 1 – Title
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG)

add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(1.2),
             "Bamboo CI/CD 順序排程實作說明", font_size=40, bold=True, color=WHITE)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(0.8),
             "四個安裝程序包 · 每天早上 9:00 · 嚴格順序執行", font_size=22, color=LIGHT)
add_text_box(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.5),
             "基於 Atlassian Bamboo · Build Plan + Stages + Cron Trigger", font_size=14, color=GRAY)

# ════════════════════════════════════════════════
# Slide 2 – 概述：問題與解法
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "問題與解決方案", font_size=32, bold=True, color=WHITE)

# left side – problem
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
         "❌ 問題",
         "• 四個安裝程序包 (1→2→3→4)\n"
         "• 必須每天早上 9:00 依序執行\n"
         "• 不可跳號、不可並行\n"
         "• 任一失敗需中止後續流程\n\n"
         "若拆成 4 個 Plan + 錯開時間 (9:00, 9:05…)\n"
         "→ 前一個超時或失敗會導致排程錯亂")
# right side – solution
add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
         "✅ 解法：單一 Plan × 四個 Stages",
         "• 一個 Build Plan 內建立 4 個 Stages\n"
         "• Stage 嚴格循序執行（Bamboo 原生保證）\n"
         "• Stage 1 成功 → Stage 2 → Stage 3 → Stage 4\n"
         "• 任一 Stage 失敗 → 自動中止，不繼續\n"
         "• 設定 Cron Trigger：0 0 9 * * ?\n\n"
         "➜ 保證順序、保證時間、失敗即停")

# ════════════════════════════════════════════════
# Slide 3 – Step 1：單一 Plan + 4 Stages 架構圖
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Step 1：建立單一 Plan + 四個 Stages", font_size=32, bold=True, color=WHITE)

y_stage = Inches(1.6)
box_w = Inches(2.2)
box_h = Inches(3.5)
gap = Inches(0.45)
start_x = Inches(0.7)

stages = [
    ("Stage 1", "安裝程序包 1\nJob: Script Task\n指令: install_pkg1.sh"),
    ("Stage 2", "安裝程序包 2\nJob: Script Task\n指令: install_pkg2.sh"),
    ("Stage 3", "安裝程序包 3\nJob: Script Task\n指令: install_pkg3.sh"),
    ("Stage 4", "安裝程序包 4\nJob: Script Task\n指令: install_pkg4.sh"),
]
colors = [ACCENT, RGBColor(0x21, 0x96, 0xF3), RGBColor(0xAB, 0x47, 0xBC), GREEN]

for i, (title, body) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_stage_box(slide, x, y_stage, box_w, box_h, title, body, color=colors[i])
    if i < 3:
        ax = x + box_w + Inches(0.05)
        add_arrow(slide, ax, y_stage + box_h // 2 - Inches(0.15), width=Inches(0.35), height=Inches(0.3))

# Plan wrapper
add_shape(slide, Inches(0.4), y_stage - Inches(0.4),
          Inches(11.5), box_h + Inches(1.0),
          fill_color=None, line_color=YELLOW, line_width=2)
add_text_box(slide, Inches(0.6), y_stage - Inches(0.35), Inches(4), Inches(0.35),
             "Build Plan: 每日順序安裝", font_size=16, bold=True, color=YELLOW)

# note
add_text_box(slide, Inches(0.8), y_stage + box_h + Inches(0.5), Inches(11), Inches(0.5),
             "💡 每個 Stage 下各有一個 Job，Job 內新增 Script Task 執行對應的安裝指令",
             font_size=14, color=LIGHT)

# ════════════════════════════════════════════════
# Slide 4 – Step 1：Bamboo UI 操作說明
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Step 1：Bamboo UI 操作步驟", font_size=32, bold=True, color=WHITE)

steps = [
    ("1", "建立 Plan", "Bamboo 首頁 → Create → Create a new plan\n輸入 Plan Name / Plan Key / Project"),
    ("2", "新增 Stage", "Plan 設定頁 → Stages 頁籤 → Add stage\n命名為「安裝程序包 1」"),
    ("3", "新增 Job & Task", "點選 Stage 下的 Default Job → Add task\n選擇 Script → 輸入安裝指令"),
    ("4", "重複 Stage 2~4", "依同樣方式建立 Stage 2 / 3 / 4\n每個 Stage 放對應的安裝 Job"),
]

for i, (num, title, body) in enumerate(steps):
    x = Inches(0.6) + (i % 2) * Inches(6.2)
    y = Inches(1.3) + (i // 2) * Inches(2.8)
    # number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    # card
    add_card(slide, x + Inches(0.6), y, Inches(5.2), Inches(2.2), title, body)

# ════════════════════════════════════════════════
# Slide 5 – Step 2：Cron 排程設定
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Step 2：設定每天 9:00 Cron 觸發", font_size=32, bold=True, color=WHITE)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
         "設定路徑",
         "Plan 編輯頁面 → Triggers 頁籤\n"
         "→ Add trigger → 選擇 Scheduled\n"
         "→ Trigger type: Cron expression")
add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(2.5),
         "Cron 運算式",
         "0 0 9 * * ?\n\n"
         "┌ 秒 (0)\n"
         "│ ┌ 分 (0)\n"
         "│ │ ┌ 時 (9)\n"
         "│ │ │ ┌ 日 (*)\n"
         "│ │ │ │ ┌ 月 (*)\n"
         "│ │ │ │ │ ┌ 星期 (?)\n"
         "0 0 9 * * ?\n\n"
         "→ 每天早上 9:00:00 觸發")

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.8),
             "💡 若需要時區調整：Bamboo 預設使用伺服器時區，請確認 Server Timezone 與目標時區一致\n"
             "可在 Administration → System Settings 中查看與設定",
             font_size=14, color=LIGHT)

# ════════════════════════════════════════════════
# Slide 6 – Step 3：錯誤處理
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Step 3：錯誤處理策略", font_size=32, bold=True, color=WHITE)

add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(3.0),
         "🛑 失敗即停（預設．建議）",
         "Stage 失敗 → 整條 Pipeline 中止\n"
         "後續 Stage 不執行\n\n"
         "適合情境：安裝包有嚴格順序\n"
         "依賴關係，前一個失敗後\n"
         "續裝無意義",
         title_color=RED)
add_card(slide, Inches(5.0), Inches(1.3), Inches(3.7), Inches(3.0),
         "⚠️ 失敗仍繼續",
         "勾選 Stage 設定中的\n"
         "「Manual Stage」\n\n"
         "或將所有安裝指令放在\n"
         "同一個 Job 的不同 Tasks\n"
         "並關閉「Fail build on task failure」",
         title_color=ORANGE)
add_card(slide, Inches(9.2), Inches(1.3), Inches(3.5), Inches(3.0),
         "📬 通知設定",
         "Bamboo → Plan 設定\n"
         "→ Notifications 頁籤\n\n"
         "可設定當 Build Failed 時\n"
         "寄送 Email / Slack 通知\n"
         "給相關負責人",
         title_color=ACCENT2)

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(1.0),
             "💡 實務建議：安裝程序包有順序相依性時，使用預設「失敗即停」即可\n"
             "         若需全部嘗試安裝（不中斷），可將四個腳本寫入同一個 Job 的不同 Task 並關閉「Fail build on task failure」",
             font_size=14, color=LIGHT)

# ════════════════════════════════════════════════
# Slide 7 – 為什麼不拆成 4 個 Plan
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "為何不使用四個獨立的 Plan？", font_size=32, bold=True, color=WHITE)

# Bad approach
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.0),
         "❌ 四個獨立 Plan + 錯開時間",
         "Plan 1 @ 9:00 → 安裝程序包 1\n"
         "Plan 2 @ 9:05 → 安裝程序包 2\n"
         "Plan 3 @ 9:10 → 安裝程序包 3\n"
         "Plan 4 @ 9:15 → 安裝程序包 4\n\n"
         "風險：\n"
         "• 若 Plan 1 執行超過 5 分鐘→Plan 2 與\n"
         "  Plan 1 同時執行，順序被打亂\n"
         "• Plan 1 失敗 → Plan 2 / 3 / 4 照常執行\n"
         "  → 安裝錯誤版本\n"
         "• 難以追蹤整體狀態（需看 4 個 Build）",
         title_color=RED)

# Good approach
add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.0),
         "✅ 一個 Plan + 四個 Stages",
         "一個 Plan，一個 Trigger\n"
         "Stage 1 → Stage 2 → Stage 3 → Stage 4\n\n"
         "優點：\n"
         "• 嚴格循序：Stage 1 完成才開始 Stage 2\n"
         "• 失敗即停：Stage 1 失敗，後面全不執行\n"
         "• 單一入口：一個 Build Number 追蹤到底\n"
         "• 狀態清晰：一頁看到完整 Pipeline 進度\n"
         "• 維護簡單：只需管理一個 Plan",
         title_color=GREEN)

# ════════════════════════════════════════════════
# Slide 8 – Summary / 總結
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), H, fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "總結：設定檢查清單", font_size=32, bold=True, color=WHITE)

checklist = [
    ("☐", "建立一個 Build Plan", "命名如「Daily Sequential Installation」"),
    ("☐", "新增 4 個 Stages", "Stage 1~4，各對應一個安裝程序包"),
    ("☐", "每個 Stage 建立一個 Job", "Job 內新增 Script Task，寫入安裝指令"),
    ("☐", "設定 Cron Trigger", "Cron expression: 0 0 9 * * ?"),
    ("☐", "確認錯誤處理方式", "建議預設失敗即停"),
    ("☐", "設定通知（選用）", "Email / Slack 通知 Build 結果"),
]

for i, (icon, title, desc) in enumerate(checklist):
    y = Inches(1.3) + i * Inches(0.85)
    add_text_box(slide, Inches(1.0), y, Inches(0.4), Inches(0.5), icon, font_size=20, color=ACCENT2)
    add_text_box(slide, Inches(1.5), y, Inches(4.5), Inches(0.5), title, font_size=18, bold=True, color=WHITE)
    add_text_box(slide, Inches(6.5), y, Inches(6), Inches(0.5), desc, font_size=16, color=LIGHT)

add_text_box(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             "💡 架構摘要：1 Plan → 4 Stages (循序) → 1 Cron Trigger (0 0 9 * * ?) → 失敗即停",
             font_size=16, bold=True, color=YELLOW)

# ── Save ──
output_path = "/Users/wesley/Downloads/lewislunora/xiangchuan/Bamboo_CI_CD_排程實作說明.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
